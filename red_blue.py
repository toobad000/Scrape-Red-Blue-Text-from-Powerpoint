#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
scrape_ppt_colors: A utility to scrape red and blue colored text from a PowerPoint file and save it to a text file.

Usage:
    python scrape_ppt_colors.py --input example.pptx

Author: Sam Hartin
Date: 2025-03-12
"""

import os
import sys
import argparse
import logging
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


def parse_arguments():
    """Parse the command-line arguments."""
    parser = argparse.ArgumentParser(description="Scrape red and blue colored text from a PowerPoint file.")
    parser.add_argument("--input", required=True, help="Input PowerPoint file.")
    return parser.parse_args()


def is_red_or_blue(color):
    """Check if the color is red or blue."""
    # Define RGB values for red and blue
    red = RGBColor(255, 0, 0)
    blue = RGBColor(0, 0, 255)

    # Check if the color is an RGBColor object
    if hasattr(color, 'rgb'):
        return color.rgb == red or color.rgb == blue
    # Check if the color is a theme color and convert it to RGB
    elif hasattr(color, 'theme_color'):
        if color.theme_color == MSO_THEME_COLOR.ACCENT_1:  # Example: Accent 1 (often red)
            return True
        elif color.theme_color == MSO_THEME_COLOR.ACCENT_2:  # Example: Accent 2 (often blue)
            return True
        else:
            return False
    else:
        return False


def scrape_colored_text(ppt_file):
    """Scrape red and blue colored text from the PowerPoint file."""
    # Load the PowerPoint presentation
    presentation = Presentation(ppt_file)
    
    # Initialize a dictionary to store the scraped text by slide number
    slide_text = {}

    # Iterate through each slide in the presentation
    for slide_number, slide in enumerate(presentation.slides, start=1):
        # Initialize a list to store the text for this slide
        slide_text[slide_number] = []
        
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if hasattr(shape, "text"):
                # Iterate through each paragraph in the shape
                for paragraph in shape.text_frame.paragraphs:
                    # Iterate through each run in the paragraph
                    for run in paragraph.runs:
                        # Check if the run has font color and if it is red or blue
                        if run.font.color and is_red_or_blue(run.font.color):
                            # Append the text to the list for this slide
                            slide_text[slide_number].append(run.text)
    
    return slide_text


def generate_output_file_name(input_file):
    """Generate the output file name based on the input file name."""
    # Get the base name of the input file (without extension)
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    # Append "_redblue.txt" to the base name
    return f"{base_name}_redblue.txt"


def main():
    """Main function for the scrape_ppt_colors script."""
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

    # Parse arguments
    args = parse_arguments()

    # Check if the input file exists
    if not os.path.isfile(args.input):
        logging.error(f"Input file {args.input} does not exist.")
        sys.exit(1)

    # Generate the output file name
    output_file = generate_output_file_name(args.input)

    # Scrape the red and blue colored text
    logging.info(f"Scraping red and blue colored text from {args.input}...")
    slide_text = scrape_colored_text(args.input)

    # Write the scraped text to the output file
    with open(output_file, 'w', encoding='utf-8') as file:
        file.write("Red and Blue Colored Text:\n\n")
        for slide_number, text_list in slide_text.items():
            if text_list:  # Only include slides with red/blue text
                file.write(f"Slide {slide_number}: {' '.join(text_list)}\n\n")
        logging.info(f"Output saved to {output_file}")

    logging.info("Scraping completed.")


if __name__ == "__main__":
    main()

#Command to Run: 
# python "c:\Users\User1\Documents\Python Scripts\red_blue.py" --input "C:\Users\User1\Documents\Python Scripts\Module7.pptx"